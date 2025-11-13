import os
import json
import mimetypes
import time
import glob
from pathlib import Path
from collections import Counter
from typing import List


def _extract_pdf(path: str) -> str:
    try:
        from PyPDF2 import PdfReader
    except Exception:
        return ""
    text = []
    try:
        with open(path, "rb") as fh:
            reader = PdfReader(fh)
            for p in reader.pages:
                try:
                    text.append(p.extract_text() or "")
                except Exception:
                    text.append("")
    except Exception:
        return ""
    return "\n".join(text).strip()


def _extract_docx(path: str) -> str:
    try:
        import docx2txt
    except Exception:
        return ""
    try:
        # docx2txt handles files itself
        return docx2txt.process(path) or ""
    except Exception:
        return ""


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""
    text = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    text.append(shp.text or "")
    except Exception:
        return ""
    return "\n".join(text).strip()


def _extract_xlsx(path: str) -> str:
    try:
        import pandas as pd
    except Exception:
        return ""
    try:
        df = pd.read_excel(path, engine="openpyxl")
        text_cols = [c for c in df.columns if df[c].dtype == object]
        if not text_cols:
            return df.to_csv(index=False)
        rows = df[text_cols].astype(str).agg(" | ".join, axis=1)
        return "\n".join(rows.tolist())
    except Exception:
        return ""


def _extract_csv(path: str) -> str:
    try:
        import pandas as pd
    except Exception:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""
    try:
        df = pd.read_csv(path)
        text_cols = [c for c in df.columns if df[c].dtype == object]
        if not text_cols:
            return df.to_csv(index=False)
        rows = df[text_cols].astype(str).agg(" | ".join, axis=1)
        return "\n".join(rows.tolist())
    except Exception:
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""


def _extract_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        try:
            with open(path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception:
            return ""


def _extract_json(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            j = json.load(f)
        return json.dumps(j, ensure_ascii=False, indent=2)
    except Exception:
        return ""


def _extract_image(path: str) -> str:
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return ""
    try:
        return pytesseract.image_to_string(Image.open(path)) or ""
    except Exception:
        return ""


def extract_top_keywords(text: str, top_k: int = 6) -> List[str]:
    if not text or len(text.split()) < 5:
        return []
    try:
        from sklearn.feature_extraction.text import TfidfVectorizer
        v = TfidfVectorizer(stop_words="english", max_features=2000)
        X = v.fit_transform([text])
        scores = X.toarray()[0]
        idxs = scores.argsort()[::-1][:top_k]
        features = v.get_feature_names_out()
        return [features[i] for i in idxs if scores[i] > 0][:top_k]
    except Exception:
        words = [w.lower().strip(".,:;()[]{}\"'") for w in text.split() if len(w) > 2]
        c = Counter(words)
        return [w for w, _ in c.most_common(top_k)]


def file_metadata(path: str) -> dict:
    p = Path(path)
    try:
        st = p.stat()
        size_bytes = st.st_size
        size_kb = round(size_bytes / 1024.0, 2)
    except Exception:
        size_bytes = None
        size_kb = "N/A"
    return {
        "filename": p.name,
        "filepath": str(p.resolve()),
        "folder": str(p.parent.resolve()),
        "size_bytes": size_bytes,
        "size_kb": size_kb,
        "modified_at": time.ctime(p.stat().st_mtime) if p.exists() else "N/A",
        "created_at": time.ctime(p.stat().st_ctime) if p.exists() else "N/A",
        "ext": p.suffix.lower(),
        "mime": mimetypes.guess_type(str(path))[0] or "unknown"
    }


def extract_single(path: str) -> dict:
    ext = Path(path).suffix.lower()
    extractor_map = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx, ".doc": _extract_docx,
        ".pptx": _extract_pptx, ".ppt": _extract_pptx,
        ".xlsx": _extract_xlsx, ".xls": _extract_xlsx,
        ".csv": _extract_csv,
        ".txt": _extract_txt, ".md": _extract_txt, ".rst": _extract_txt,
        ".json": _extract_json,
        ".png": _extract_image, ".jpg": _extract_image, ".jpeg": _extract_image,
        ".bmp": _extract_image, ".tiff": _extract_image
    }

    extractor = extractor_map.get(ext, _extract_txt)
    text = extractor(path) or ""
    meta = file_metadata(path)

    word_count = len(text.split())
    preview = (text[:300] + "...") if len(text) > 300 else text
    keywords = extract_top_keywords(text, top_k=6)

    return {
        "id": str(abs(hash(meta["filepath"])) % (10**9)),
        "title": meta["filename"],
        "content": text,
        "preview": preview,
        "metadata": {**meta, "word_count": word_count, "top_keywords": keywords}
    }


def _normalize_inputs(inputs: List[str]) -> List[str]:
    paths = []
    for inp in inputs:
        if ";" in inp:
            for part in inp.split(";"):
                paths.append(part.strip())
        else:
            paths.append(inp)

    resolved = []
    for p in paths:
        matches = glob.glob(os.path.expanduser(p), recursive=True)
        if matches:
            resolved.extend(matches)
        else:
            resolved.append(p)

    unique = []
    seen = set()
    for p in resolved:
        try:
            abs_p = str(Path(p).resolve())
        except Exception:
            abs_p = os.path.abspath(p)
        if abs_p not in seen:
            seen.add(abs_p)
            unique.append(abs_p)
    return unique


def extract_documents_from_paths(inputs: List[str]) -> List[dict]:
    all_docs = []
    inputs = _normalize_inputs(inputs)

    for p in inputs:
        path = Path(p)
        if not path.exists():
            print(f" Skipping missing path: {p}")
            continue

        if path.is_dir():
            for root, _, files in os.walk(path):
                for fn in files:
                    file_path = os.path.join(root, fn)
                    print(f"  - Extracting: {fn}")
                    try:
                        doc = extract_single(file_path)
                        if doc["content"] and len(doc["content"].strip()) > 10:
                            all_docs.append(doc)
                        else:
                            print("    ⚠️ no textual content or too small; skipping")
                    except Exception as e:
                        print("    ❌ extraction error:", e)
        else:
            print(f"  - Extracting single file: {p}")
            try:
                doc = extract_single(str(p))
                if doc["content"] and len(doc["content"].strip()) > 10:
                    all_docs.append(doc)
                else:
                    print("    ⚠️ no textual content or too small; skipping")
            except Exception as e:
                print("    ❌ extraction error:", e)

    print(f"\n Extracted {len(all_docs)} document(s).")
    return all_docs


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python document_extractor.py <folder_or_file> [more_paths...]")
        sys.exit(1)

    docs = extract_documents_from_paths(sys.argv[1:])
    output = "extracted_documents.json"
    with open(output, "w", encoding="utf-8") as f:
        json.dump(docs, f, indent=2, ensure_ascii=False)
    print(f" Saved {len(docs)} docs to {output}")
